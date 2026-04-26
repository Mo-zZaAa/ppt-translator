#!/usr/bin/env python3
"""
PPT Translator — Korean → English
Uses Upstage Solar mini API for batch paragraph-level translation.
Preserves all formatting (font, size, color, position, images).

Key design:
- Translate at PARAGRAPH level (not run level) to preserve sentence context
- First run of each paragraph gets the full translated text
- Remaining runs are cleared (format attributes preserved)
- All translation done via Upstage Solar API — zero Claude tokens used
"""

import json
import re
import sys
import argparse
import requests
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

UPSTAGE_API_URL = "https://api.upstage.ai/v1/solar/chat/completions"
KOREAN_RE = re.compile(r"[가-힣ㄱ-ㅎㅏ-ㅣ]")


def has_korean(text: str) -> bool:
    return bool(KOREAN_RE.search(text))


# ---------------------------------------------------------------------------
# Collection
# ---------------------------------------------------------------------------

def collect_paragraphs(prs):
    """
    Walk all shapes and collect paragraphs that contain Korean text.
    Returns list of (para_obj, full_text) tuples.
    full_text = concatenation of all run texts in the paragraph.
    """
    results = []

    def _from_tf(tf):
        for para in tf.paragraphs:
            full_text = "".join(r.text for r in para.runs)
            if has_korean(full_text):
                results.append((para, full_text))

    for slide in prs.slides:
        for shape in slide.shapes:
            _process_shape(shape, _from_tf)

    return results


def _process_shape(shape, callback):
    """Recursively process shape, handling text frames, tables, groups."""
    try:
        if shape.has_text_frame:
            callback(shape.text_frame)
    except Exception:
        pass

    try:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    callback(cell.text_frame)
    except Exception:
        pass

    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                _process_shape(child, callback)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Translation
# ---------------------------------------------------------------------------

def _translate_chunk(chunk: list[str], api_key: str) -> dict[str, str]:
    """Send one chunk to Upstage Solar mini and return {original: translated}."""
    system_prompt = (
        "You are a professional Korean-to-English translator specializing in "
        "business consulting presentations. "
        "The user sends a JSON array of Korean text strings "
        "(each string is one paragraph from a slide). "
        "Return a JSON object where each key is the EXACT original Korean string "
        "and the value is the English translation. "
        "Rules:\n"
        "- Keep proper nouns, brand names, company names, and numbers unchanged.\n"
        "- Translate naturally and concisely — match the original brevity.\n"
        "- Do NOT add explanations or parentheses.\n"
        "- Return ONLY the JSON object, no markdown, no extra text."
    )
    payload = {
        "model": "solar-mini",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": json.dumps(chunk, ensure_ascii=False)},
        ],
        "response_format": {"type": "json_object"},
        "temperature": 0.1,
        "max_tokens": 4096,
    }
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    resp = requests.post(UPSTAGE_API_URL, json=payload, headers=headers, timeout=60)
    resp.raise_for_status()
    content = resp.json()["choices"][0]["message"]["content"]

    # Sanitize: remove unescaped control characters that break JSON parsing
    # Replace literal newlines/tabs inside JSON string values with escaped versions
    content = re.sub(r'(?<!\\)\n', ' ', content)
    content = re.sub(r'(?<!\\)\r', ' ', content)
    content = re.sub(r'(?<!\\)\t', ' ', content)
    # Remove other ASCII control chars (0x00-0x1f) except escaped ones
    content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', content)

    return json.loads(content)


def batch_translate(texts: list[str], api_key: str, chunk_size: int = 20) -> dict[str, str]:
    """
    Translate all unique Korean strings via Upstage Solar.
    Chunks to avoid token/timeout limits. Retries once on timeout.
    Returns {original_korean: english_translation}.
    """
    unique = list(dict.fromkeys(texts))
    result = {}
    chunks = [unique[i:i + chunk_size] for i in range(0, len(unique), chunk_size)]
    print(f"Translating {len(unique)} unique paragraphs in {len(chunks)} chunks via Upstage Solar...")

    for i, chunk in enumerate(chunks, 1):
        print(f"  Chunk {i}/{len(chunks)} ({len(chunk)} strings)...", end=" ", flush=True)
        success = False
        for attempt in range(2):  # retry once on failure
            try:
                translated = _translate_chunk(chunk, api_key)
                result.update(translated)
                print("✓")
                success = True
                break
            except requests.exceptions.Timeout:
                if attempt == 0:
                    print(f"timeout, retrying...", end=" ", flush=True)
                else:
                    print("✗ timeout again, skipping chunk")
            except Exception as e:
                print(f"✗ ERROR: {e}")
                break
        if not success:
            # fallback: keep originals for failed chunk
            for t in chunk:
                result[t] = t

    return result


# ---------------------------------------------------------------------------
# Application
# ---------------------------------------------------------------------------

def apply_translation(para, translated_text: str):
    """
    Put translated_text into first run, clear the rest.
    Preserves run formatting attributes (font, size, color, bold, etc.)
    """
    runs = para.runs
    if not runs:
        return
    runs[0].text = translated_text
    for run in runs[1:]:
        run.text = ""


# ---------------------------------------------------------------------------
# Main pipeline
# ---------------------------------------------------------------------------

def translate_pptx(input_path: str, output_path: str, api_key: str) -> int:
    """
    Full translation pipeline.
    Returns number of paragraphs translated.
    """
    print(f"Loading: {input_path}")
    prs = Presentation(input_path)

    # Step 1: Collect Korean paragraphs
    paragraphs = collect_paragraphs(prs)
    if not paragraphs:
        print("No Korean text found. Saving as-is.")
        prs.save(output_path)
        return 0

    print(f"Found {len(paragraphs)} Korean paragraphs.")

    # Step 2: Batch translate via Upstage Solar
    texts = [text for _, text in paragraphs]
    translation_map = batch_translate(texts, api_key)

    # Step 3: Apply translations
    translated_count = 0
    for para, original_text in paragraphs:
        translated = translation_map.get(original_text)
        if translated and translated != original_text:
            apply_translation(para, translated)
            translated_count += 1
        elif not translated:
            # Try with stripped text
            stripped = original_text.strip()
            translated = translation_map.get(stripped)
            if translated and translated != original_text:
                apply_translation(para, translated)
                translated_count += 1

    # Step 4: Save
    prs.save(output_path)
    print(f"\n✅ Done: {translated_count}/{len(paragraphs)} paragraphs translated.")
    print(f"   Saved → {output_path}")
    return translated_count


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Translate Korean PPT to English using Upstage Solar API"
    )
    parser.add_argument("input", help="Input .pptx file path")
    parser.add_argument(
        "-o", "--output",
        help="Output file path (default: <input>_EN.pptx)"
    )
    parser.add_argument(
        "--api-key",
        required=True,
        help="Upstage API key (https://console.upstage.ai)"
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=40,
        help="Number of paragraphs per API call (default: 40)"
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = args.output or str(input_path.with_stem(input_path.stem + "_EN"))
    translate_pptx(str(input_path), output_path, args.api_key)


if __name__ == "__main__":
    main()
